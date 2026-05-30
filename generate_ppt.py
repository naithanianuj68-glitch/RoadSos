import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Custom Palette
    BG_COLOR = RGBColor(11, 15, 25)          # Dark Navy (#0B0F19)
    TEXT_WHITE = RGBColor(255, 255, 255)     # Clean White
    TEXT_MUTED = RGBColor(148, 163, 184)     # Slate Grey (#94A3B8)
    COLOR_RED = RGBColor(239, 68, 68)        # Emergency Red (#EF4444)
    COLOR_GREEN = RGBColor(16, 185, 129)     # Safety Green (#10B981)
    COLOR_BLUE = RGBColor(59, 130, 246)      # Tech Blue (#3B82F6)
    COLOR_PURPLE = RGBColor(139, 92, 246)    # Tech Purple (#8B5CF6)
    COLOR_CARD = RGBColor(19, 28, 48)        # Card Background (#131C30)
    COLOR_CARD_BORDER = RGBColor(38, 50, 77) # Card Border (#26324D)
    
    FONT_FAMILY = 'Trebuchet MS'  # Modern geometric sans-serif
    
    blank_slide_layout = prs.slide_layouts[6] # Blank layout

    # Helper to set solid background color
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper to draw cards
    def draw_card(slide, left, top, width, height, border_color=COLOR_CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    # Helper to draw smartphone mockup
    def draw_phone_mockup(slide, left, top, width, height, image_path, border_color):
        # 1. Outer phone frame (chassis bezel)
        chassis = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        chassis.fill.solid()
        chassis.fill.fore_color.rgb = RGBColor(15, 15, 20) # Dark grey chassis
        chassis.line.color.rgb = border_color
        chassis.line.width = Pt(2.5)
        
        # 2. Speaker notch at top
        notch_w = Inches(0.6)
        notch_h = Inches(0.08)
        notch_l = left + (width - notch_w) / 2
        notch_t = top + Inches(0.12)
        notch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, notch_l, notch_t, notch_w, notch_h)
        notch.fill.solid()
        notch.fill.fore_color.rgb = RGBColor(40, 40, 45)
        notch.line.fill.background()
        
        # 3. Home indicator bar at bottom
        bar_w = Inches(0.8)
        bar_h = Inches(0.04)
        bar_l = left + (width - bar_w) / 2
        bar_t = top + height - Inches(0.15)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_l, bar_t, bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(120, 120, 125)
        bar.line.fill.background()
        
        # 4. Screenshot Image inset
        inset_x = Inches(0.08)
        inset_y_top = Inches(0.3)
        inset_y_bottom = Inches(0.3)
        img_w = width - (inset_x * 2)
        img_h = height - inset_y_top - inset_y_bottom
        img_l = left + inset_x
        img_t = top + inset_y_top
        
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, img_l, img_t, width=img_w, height=img_h)

    # Helper to add formatted titles & divider lines
    def add_slide_header(slide, tag_text, title_text, color_tag=COLOR_RED):
        # Tag pill shape
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(1.8), Inches(0.35))
        pill.fill.solid()
        pill.fill.fore_color.rgb = color_tag
        pill.line.fill.background()
        tf_pill = pill.text_frame
        tf_pill.word_wrap = True
        p_pill = tf_pill.paragraphs[0]
        p_pill.text = tag_text.upper()
        p_pill.alignment = PP_ALIGN.CENTER
        p_pill.font.name = FONT_FAMILY
        p_pill.font.size = Pt(11)
        p_pill.font.bold = True
        p_pill.font.color.rgb = TEXT_WHITE
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.7))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = FONT_FAMILY
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = TEXT_WHITE
        
        # Accent Line Divider
        connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.03))
        connector.fill.solid()
        connector.fill.fore_color.rgb = color_tag
        connector.line.fill.background()

    # ==================== SLIDE 1: WELCOME ====================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide1)
    
    # Large Decorative Background Glow Shape
    glow = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.666), Inches(0.5), Inches(6.0), Inches(6.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(16, 20, 36)
    glow.line.fill.background()
    
    # Insert Logo
    logo_path = 'public/logo.png'
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(5.666), Inches(1.3), width=Inches(2.0))
        
    # App Title
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "RoadSoS"
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_FAMILY
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    
    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.333), Inches(0.8))
    tf2 = sub_box.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "The Ultimate Offline-First Emergency Network"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT_FAMILY
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_WHITE
    
    p2_sub = tf2.add_paragraph()
    p2_sub.text = "EMPOWERING ACCIDENT VICTIMS & BYSTANDERS DURING THE GOLDEN HOUR"
    p2_sub.alignment = PP_ALIGN.CENTER
    p2_sub.font.name = FONT_FAMILY
    p2_sub.font.size = Pt(11)
    p2_sub.font.bold = True
    p2_sub.font.color.rgb = TEXT_MUTED
    p2_sub.space_before = Pt(8)

    # Footer/Hackathon info
    team_box = slide1.shapes.add_textbox(Inches(1.0), Inches(6.1), Inches(11.333), Inches(0.6))
    tf3 = team_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "NATIONAL ROAD SAFETY HACKATHON 2026  |  TEAM ROADSOS INNOVATORS"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.name = FONT_FAMILY
    p3.font.size = Pt(11)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_RED

    # ==================== SLIDE 2: THE PROBLEM ====================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide2)
    add_slide_header(slide2, "The Problem", "Golden Hour Delays Cost Lives on Highways", COLOR_RED)
    
    # Left hook description
    desc_box = slide2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "Highway emergency response rates are critical. Survival rates drop by 60% for every 10-minute delay during the 'Golden Hour' (first 60 mins):"
    p_desc.font.name = FONT_FAMILY
    p_desc.font.size = Pt(22)
    p_desc.font.color.rgb = TEXT_WHITE
    p_desc.space_after = Pt(20)
    
    p_subdesc = tf_desc.add_paragraph()
    p_subdesc.text = "Traditional calling applications, medical guidance databases, and maps depend heavily on stable high-speed cellular networks. On critical highway segments, this reliance introduces single-points-of-failure."
    p_subdesc.font.name = FONT_FAMILY
    p_subdesc.font.size = Pt(14)
    p_subdesc.font.color.rgb = TEXT_MUTED
    p_subdesc.space_before = Pt(10)

    # Right Cards
    draw_card(slide2, Inches(6.6), Inches(2.0), Inches(5.9), Inches(2.1), COLOR_RED)
    c1_box = slide2.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.5), Inches(1.9))
    tf_c1 = c1_box.text_frame
    tf_c1.word_wrap = True
    p_c1_title = tf_c1.paragraphs[0]
    p_c1_title.text = "📶 Cellular Dead Zones (OSM Data)"
    p_c1_title.font.name = FONT_FAMILY
    p_c1_title.font.size = Pt(18)
    p_c1_title.font.bold = True
    p_c1_title.font.color.rgb = COLOR_RED
    p_c1_body = tf_c1.add_paragraph()
    p_c1_body.text = "Over 40% of national highway corridors pass through cellular coverage gaps or dead zones. Standard API-based apps (like Google Maps or Uber) stall immediately with zero bars of signal."
    p_c1_body.font.name = FONT_FAMILY
    p_c1_body.font.size = Pt(13)
    p_c1_body.font.color.rgb = TEXT_MUTED
    p_c1_body.space_before = Pt(6)

    draw_card(slide2, Inches(6.6), Inches(4.5), Inches(5.9), Inches(2.1), COLOR_RED)
    c2_box = slide2.shapes.add_textbox(Inches(6.8), Inches(4.6), Inches(5.5), Inches(1.9))
    tf_c2 = c2_box.text_frame
    tf_c2.word_wrap = True
    p_c2_title = tf_c2.paragraphs[0]
    p_c2_title.text = "🧠 Victim Trauma & Bystander Panic"
    p_c2_title.font.name = FONT_FAMILY
    p_c2_title.font.size = Pt(18)
    p_c2_title.font.bold = True
    p_c2_title.font.color.rgb = COLOR_RED
    p_c2_body = tf_c2.add_paragraph()
    p_c2_body.text = "Victims suffer physical shock, limiting phone interaction. Bystanders, driven by panic, fail to communicate GPS coordinate metrics or recall regional emergency helpline directories."
    p_c2_body.font.name = FONT_FAMILY
    p_c2_body.font.size = Pt(13)
    p_c2_body.font.color.rgb = TEXT_MUTED
    p_c2_body.space_before = Pt(6)

    # ==================== SLIDE 3: OUR SOLUTION ====================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide3)
    add_slide_header(slide3, "Our Solution", "Autonomous, Offline-First Rescue Ecosystem", COLOR_GREEN)
    
    # Left Column: Description & Bullet Cards
    sol_box = slide3.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.0), Inches(4.5))
    tf_sol = sol_box.text_frame
    tf_sol.word_wrap = True
    p_sol = tf_sol.paragraphs[0]
    p_sol.text = "RoadSoS is an offline-first PWA that bypasses cellular dependence to ensure continuous safety backup:"
    p_sol.font.name = FONT_FAMILY
    p_sol.font.size = Pt(20)
    p_sol.font.bold = True
    p_sol.font.color.rgb = TEXT_WHITE
    p_sol.space_after = Pt(16)
    
    points = [
      ("🌐 IndexedDB Spatial Caching", "Uses Dexie.js to store over 1,500+ regional hospital coordinates, pharmacy locations, and local police helpline directories on client storage. Geolocation signals map directions offline."),
      ("📱 Zero Installation Footprint", "Progressive Web App (PWA) boot speeds complete in <1.5 seconds. Instantly accessible via SMS web links or vehicle QR code stickers without app store downloads.")
    ]
    for title, desc in points:
        p_pt_t = tf_sol.add_paragraph()
        p_pt_t.text = title
        p_pt_t.font.name = FONT_FAMILY
        p_pt_t.font.size = Pt(17)
        p_pt_t.font.bold = True
        p_pt_t.font.color.rgb = COLOR_GREEN
        p_pt_t.space_before = Pt(10)
        
        p_pt_d = tf_sol.add_paragraph()
        p_pt_d.text = desc
        p_pt_d.font.name = FONT_FAMILY
        p_pt_d.font.size = Pt(13)
        p_pt_d.font.color.rgb = TEXT_MUTED
        p_pt_d.space_before = Pt(4)

    # Right Column: Screenshot container centered and aligned to phone aspect ratio
    draw_phone_mockup(slide3, Inches(8.65), Inches(1.8), Inches(2.6), Inches(5.0), 'public/app_mockup.png', COLOR_GREEN)

    # ==================== SLIDE 4: CORE FEATURES ====================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide4)
    add_slide_header(slide4, "Core Engine", "Hardware SOS Blackbox & AI Safety Copilot", COLOR_BLUE)
    
    # Left Column: Heading
    lh_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.4), Inches(0.8))
    tf_lh = lh_box.text_frame
    tf_lh.word_wrap = True
    
    p_lh = tf_lh.paragraphs[0]
    p_lh.text = "Dual-Shield Safety Protocols"
    p_lh.font.name = FONT_FAMILY
    p_lh.font.size = Pt(22)
    p_lh.font.bold = True
    p_lh.font.color.rgb = TEXT_WHITE
    
    p_lh_sub = tf_lh.add_paragraph()
    p_lh_sub.text = "Combines real-time proactive driver speed safety alerts with autonomous reactive triggers:"
    p_lh_sub.font.name = FONT_FAMILY
    p_lh_sub.font.size = Pt(13)
    p_lh_sub.font.color.rgb = TEXT_MUTED
    
    # Feature Column 1 (Left)
    col1_box = slide4.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(3.1), Inches(4.5))
    tf_col1 = col1_box.text_frame
    tf_col1.word_wrap = True
    
    features_left = [
      ("🔊 Audio Siren wave", "Generates raw 440Hz sawtooth frequencies using Web Audio API to alert nearby rescuers in dark conditions."),
      ("🎙️ Forensic Recorder", "Local MediaRecorder saves 15 seconds of accident audio logs directly to browser indexed storage."),
      ("📳 Impact Shake Detect", "Uses device accelerometer filters to detect structural G-force and auto-dial 112 emergency numbers.")
    ]
    
    first = True
    for title, desc in features_left:
        p_title = tf_col1.paragraphs[0] if first else tf_col1.add_paragraph()
        first = False
        p_title.text = title
        p_title.font.name = FONT_FAMILY
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_BLUE
        if p_title != tf_col1.paragraphs[0]:
            p_title.space_before = Pt(10)
            
        p_desc = tf_col1.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_FAMILY
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(2)
        
    # Feature Column 2 (Right)
    col2_box = slide4.shapes.add_textbox(Inches(4.1), Inches(2.6), Inches(3.1), Inches(4.5))
    tf_col2 = col2_box.text_frame
    tf_col2.word_wrap = True
    
    features_right = [
      ("⚡ GPS Speed Monitor", "Watches speed velocity via Geolocation API, flashing warning lights and alert buzzers to check limits."),
      ("📝 Report Generator", "Instantly packs lat/lng coordinates, address descriptors, and audio links into text reports."),
      ("📍 Hotspot Map Alerts", "Pre-loads accident-prone blackspot coordinates, warning driver upon zone entry offline.")
    ]
    
    first = True
    for title, desc in features_right:
        p_title = tf_col2.paragraphs[0] if first else tf_col2.add_paragraph()
        first = False
        p_title.text = title
        p_title.font.name = FONT_FAMILY
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_BLUE
        if p_title != tf_col2.paragraphs[0]:
            p_title.space_before = Pt(10)
            
        p_desc = tf_col2.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_FAMILY
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.space_before = Pt(2)

    # Right Column: Screenshot container centered and aligned to phone aspect ratio
    draw_phone_mockup(slide4, Inches(8.65), Inches(1.8), Inches(2.6), Inches(5.0), 'public/ai_mockup.png', COLOR_BLUE)

    # ==================== SLIDE 5: SYSTEM ARCHITECTURE ====================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide5)
    add_slide_header(slide5, "Technical Stack", "Resilient System Architecture Layers", COLOR_PURPLE)
    
    # 3 Column explanation blocks as cards (optimised spacing & visual indicators)
    card_w = Inches(3.6)
    card_h = Inches(4.6)
    positions = [Inches(0.8), Inches(4.85), Inches(8.9)]
    themes = [COLOR_PURPLE, COLOR_GREEN, COLOR_BLUE]
    
    # Card 1 Contents
    bullets1 = [
        "React + Vite Framework: Single Page Application layout boots in milliseconds.",
        "Tailwind CSS: Dark-mode interface optimized for outdoor high-glare and night emergency visibility.",
        "Leaflet.js Mapping: Implements canvas-based offline tile render hooks."
    ]
    # Card 2 Contents
    bullets2 = [
        "PWA Service Worker: Caches main application shell and scripts via Workbox.",
        "IndexedDB Database: Handles regional hospital listings, police caches, and map files local queries.",
        "Local Rule-Engine: Chatbot runs on-device rule matching for first-aid without internet."
    ]
    # Card 3 Contents
    bullets3 = [
        "OpenStreetMap API: Fetches raw public spatial data for emergency services.",
        "Claude LLM: Parses complex symptoms and multilingual questions when online.",
        "Hardware Sensor APIs: Binds Accelerometer, AudioContext, and Media Recording tools."
    ]
    
    all_bullets = [bullets1, bullets2, bullets3]
    titles = ["1. Client Frontend", "2. Offline Guard", "3. Integration APIs"]
    
    for i, left in enumerate(positions):
        color = themes[i]
        # Draw background card
        draw_card(slide5, left, Inches(2.0), card_w, card_h, color)
        
        # Decorative top bar
        top_bar = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), card_w, Inches(0.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()
        
        # Header TextBox
        h_box = slide5.shapes.add_textbox(left + Inches(0.25), Inches(2.2), card_w - Inches(0.5), Inches(0.55))
        tf_h = h_box.text_frame
        tf_h.word_wrap = True
        tf_h.margin_left = tf_h.margin_top = tf_h.margin_right = tf_h.margin_bottom = 0
        p_t = tf_h.paragraphs[0]
        p_t.text = titles[i]
        p_t.font.name = FONT_FAMILY
        p_t.font.size = Pt(17.5)
        p_t.font.bold = True
        p_t.font.color.rgb = color
        
        # Accent separator line under header
        sep = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.25), Inches(2.8), card_w - Inches(0.5), Inches(0.02))
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(38, 50, 77) # Muted card border color
        sep.line.fill.background()
        
        # Body TextBox
        b_box = slide5.shapes.add_textbox(left + Inches(0.25), Inches(2.95), card_w - Inches(0.5), Inches(3.5))
        tf_b = b_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        
        first = True
        for bullet in all_bullets[i]:
            p = tf_b.paragraphs[0] if first else tf_b.add_paragraph()
            first = False
            p.text = "• " + bullet
            p.font.name = FONT_FAMILY
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED
            p.space_before = Pt(6)

    # ==================== SLIDE 6: THE IMPACT ====================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide6)
    add_slide_header(slide6, "The Impact", "Scalable, Zero Cost, Lifesaving Platform", COLOR_GREEN)
    
    # Visual grid of cards for impact stats
    imp_positions = [Inches(0.8), Inches(4.85), Inches(8.9)]
    imp_themes = [COLOR_GREEN, COLOR_BLUE, COLOR_RED]
    imp_metrics = ["$0", "100%", "Instant"]
    imp_labels = ["Map API Licensing Fees", "Universal Location Coding", "Offline Zero-Bar Access"]
    imp_descs = [
        "Commercial maps cost $7 per 1,000 lookups. RoadSoS leverages free OpenStreetMap datasets and local cache indexing. This design eliminates runtime overheads, allowing global scaling at zero platform fees.",
        "OSM maps cover all roads, toll plazas, and emergency centers globally. The app matches coordinate boundaries locally in any country (e.g. India, EU, US), updating emergency contact numbers dynamically.",
        "In trauma situations, installation delays lead to negative outcomes. PWA formats boot directly in mobile browsers. Essential siren systems, GPS alerts, and first-aid run with zero network active."
    ]
    
    for i, left in enumerate(imp_positions):
        color = imp_themes[i]
        # Draw background card
        draw_card(slide6, left, Inches(2.0), card_w, card_h, color)
        
        # Decorative top bar
        top_bar = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), card_w, Inches(0.12))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = color
        top_bar.line.fill.background()
        
        # Metric Highlight Circle Background Glow
        glow_size = Inches(0.8)
        glow_c = slide6.shapes.add_shape(MSO_SHAPE.OVAL, left + card_w - Inches(1.05), Inches(2.25), glow_size, glow_size)
        glow_c.fill.solid()
        glow_c.fill.fore_color.rgb = RGBColor(15, 23, 42)
        glow_c.line.color.rgb = color
        glow_c.line.width = Pt(1.5)
        
        # Small icon or tag number inside the glow
        tf_glow = glow_c.text_frame
        tf_glow.margin_left = tf_glow.margin_top = tf_glow.margin_right = tf_glow.margin_bottom = 0
        p_glow = tf_glow.paragraphs[0]
        p_glow.text = f"0{i+1}"
        p_glow.alignment = PP_ALIGN.CENTER
        p_glow.font.name = FONT_FAMILY
        p_glow.font.size = Pt(11)
        p_glow.font.bold = True
        p_glow.font.color.rgb = color
        
        # Metric Value TextBox
        m_box = slide6.shapes.add_textbox(left + Inches(0.25), Inches(2.25), card_w - Inches(1.3), Inches(0.85))
        tf_m = m_box.text_frame
        tf_m.word_wrap = True
        tf_m.margin_left = tf_m.margin_top = tf_m.margin_right = tf_m.margin_bottom = 0
        p_m = tf_m.paragraphs[0]
        p_m.text = imp_metrics[i]
        p_m.font.name = FONT_FAMILY
        p_m.font.size = Pt(44)
        p_m.font.bold = True
        p_m.font.color.rgb = color
        
        # Accent separator line under header
        sep = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.25), Inches(3.2), card_w - Inches(0.5), Inches(0.02))
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(38, 50, 77)
        sep.line.fill.background()
        
        # Label Title Box
        lbl_box = slide6.shapes.add_textbox(left + Inches(0.25), Inches(3.35), card_w - Inches(0.5), Inches(0.4))
        tf_l = lbl_box.text_frame
        tf_l.word_wrap = True
        tf_l.margin_left = tf_l.margin_top = tf_l.margin_right = tf_l.margin_bottom = 0
        p_l = tf_l.paragraphs[0]
        p_l.text = imp_labels[i]
        p_l.font.name = FONT_FAMILY
        p_l.font.size = Pt(13)
        p_l.font.bold = True
        p_l.font.color.rgb = TEXT_WHITE
        
        # Description TextBox
        d_box = slide6.shapes.add_textbox(left + Inches(0.25), Inches(3.8), card_w - Inches(0.5), Inches(2.6))
        tf_d = d_box.text_frame
        tf_d.word_wrap = True
        tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = 0
        p_d = tf_d.paragraphs[0]
        p_d.text = imp_descs[i]
        p_d.font.name = FONT_FAMILY
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(0)

    # ==================== SLIDE 7: THANK YOU ====================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide7)
    
    # Subtitle contrast glow
    glow7 = slide7.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.666), Inches(1.5), Inches(6.0), Inches(4.5))
    glow7.fill.solid()
    glow7.fill.fore_color.rgb = RGBColor(16, 20, 36)
    glow7.line.fill.background()
    
    # Large Thank You Text
    ty_box = slide7.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(1.2))
    tf_ty = ty_box.text_frame
    tf_ty.word_wrap = True
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "Thank You!"
    p_ty.alignment = PP_ALIGN.CENTER
    p_ty.font.name = FONT_FAMILY
    p_ty.font.size = Pt(64)
    p_ty.font.bold = True
    p_ty.font.color.rgb = TEXT_WHITE
    
    p_tys = tf_ty.add_paragraph()
    p_tys.text = "Let's make our highways safer, together."
    p_tys.alignment = PP_ALIGN.CENTER
    p_tys.font.name = FONT_FAMILY
    p_tys.font.size = Pt(20)
    p_tys.font.color.rgb = TEXT_MUTED
    p_tys.space_before = Pt(8)
    
    # Live URL Details Box
    draw_card(slide7, Inches(3.666), Inches(3.6), Inches(6.0), Inches(2.2), COLOR_RED)
    info_box = slide7.shapes.add_textbox(Inches(3.766), Inches(3.7), Inches(5.8), Inches(2.0))
    tf_info = info_box.text_frame
    tf_info.word_wrap = True
    
    p_i_t = tf_info.paragraphs[0]
    p_i_t.text = "📱 Try RoadSoS Live Demo"
    p_i_t.alignment = PP_ALIGN.CENTER
    p_i_t.font.name = FONT_FAMILY
    p_i_t.font.size = Pt(20)
    p_i_t.font.bold = True
    p_i_t.font.color.rgb = COLOR_RED
    
    p_i_b = tf_info.add_paragraph()
    p_i_b.text = "Instant loading, offline cached PWA ready for tests."
    p_i_b.alignment = PP_ALIGN.CENTER
    p_i_b.font.name = FONT_FAMILY
    p_i_b.font.size = Pt(13)
    p_i_b.font.color.rgb = TEXT_MUTED
    p_i_b.space_before = Pt(6)
    
    p_i_url = tf_info.add_paragraph()
    p_i_url.text = "LIVE REPOSITORY: github.com/TeamRoadSoS/app"
    p_i_url.alignment = PP_ALIGN.CENTER
    p_i_url.font.name = FONT_FAMILY
    p_i_url.font.size = Pt(14)
    p_i_url.font.bold = True
    p_i_url.font.color.rgb = TEXT_WHITE
    p_i_url.space_before = Pt(10)

    # Save presentation
    output_filename = "RoadSoS_Pitch_Deck.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as {output_filename}")

if __name__ == "__main__":
    create_presentation()
